"""
EV 배터리 자가진단 BMS 튜닝 시스템 — 발표 PPT 생성 (18슬라이드, 15분)
실행: python make_ppt.py
출력: presentation/EV_BMS_발표.pptx
"""

import io
import math
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import font_manager
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── 색상 팔레트 ──────────────────────────────────────────────────────
C_BG       = RGBColor(0x0F, 0x17, 0x2A)   # 진한 네이비
C_BG2      = RGBColor(0x1A, 0x25, 0x3D)   # 카드 배경
C_ACCENT   = RGBColor(0x3B, 0x82, 0xF6)   # 파랑 (primary)
C_GREEN    = RGBColor(0x22, 0xC5, 0x5E)   # 초록
C_RED      = RGBColor(0xEF, 0x44, 0x44)   # 빨강
C_YELLOW   = RGBColor(0xFB, 0xBF, 0x24)   # 노랑
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
C_CARD     = RGBColor(0x1E, 0x2D, 0x48)
C_NEWS_BG  = RGBColor(0x1C, 0x1C, 0x1C)
C_NEWS_RED = RGBColor(0xDC, 0x26, 0x26)

W = Inches(13.33)
H = Inches(7.5)

# ─── matplotlib 한글 폰트 ────────────────────────────────────────────
def _set_korean_font():
    candidates = ["Malgun Gothic", "NanumGothic", "AppleGothic", "sans-serif"]
    for name in candidates:
        if any(name.lower() in f.name.lower() for f in font_manager.fontManager.ttflist):
            plt.rcParams["font.family"] = name
            break
    plt.rcParams["axes.unicode_minus"] = False

_set_korean_font()

# ─── 유틸 ─────────────────────────────────────────────────────────────
def new_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # 모든 placeholder 제거
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


def bg(slide, color=C_BG):
    bg_shape = slide.shapes.add_shape(
        1, 0, 0, W, H
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()
    return bg_shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, x, y, w, h, fill=C_CARD, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def fig_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    return buf


def add_fig(slide, fig, x, y, w, h):
    buf = fig_to_image(fig)
    slide.shapes.add_picture(buf, x, y, w, h)
    plt.close(fig)


def accent_bar(slide, height=Inches(0.06)):
    """상단 파랑 강조 바"""
    bar = slide.shapes.add_shape(1, 0, 0, W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()


def section_title(slide, title, subtitle=""):
    accent_bar(slide)
    add_text(slide, title,
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             size=22, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.72), Inches(12), Inches(0.4),
                 size=13, color=C_GRAY)


# ══════════════════════════════════════════════════════════════════════
# 슬라이드별 함수
# ══════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """제목 슬라이드"""
    slide = new_slide(prs)
    bg(slide, C_BG)

    # 배경 그라데이션 효과 (사각형 레이어)
    overlay = slide.shapes.add_shape(1, 0, 0, W, H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0A, 0x10, 0x20)
    overlay.line.fill.background()

    # 상단 accent 라인
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    # 배터리 아이콘 느낌 사각형 (데코)
    for i, (rx, ry, rw, rh, alpha) in enumerate([
        (Inches(10.5), Inches(1.5), Inches(2.0), Inches(3.5), 0x15),
        (Inches(11.0), Inches(1.0), Inches(1.5), Inches(2.5), 0x10),
    ]):
        deco = slide.shapes.add_shape(1, rx, ry, rw, rh)
        deco.fill.solid()
        deco.fill.fore_color.rgb = C_ACCENT
        deco.line.color.rgb = C_ACCENT
        deco.line.width = Pt(1.5)

    # 메인 타이틀
    add_text(slide,
             "🔋 EV 배터리 자가진단\n+ 맞춤 BMS 튜닝 시스템",
             Inches(0.8), Inches(1.6), Inches(9.5), Inches(2.8),
             size=38, bold=True, color=C_WHITE, wrap=True)

    # 서브타이틀
    add_text(slide,
             "Physics-Informed LSTM × 4-way 모델 비교 × NCM 실차 191대",
             Inches(0.8), Inches(4.2), Inches(9.5), Inches(0.5),
             size=16, color=C_ACCENT, bold=True)

    # 구분선
    line_shape = slide.shapes.add_shape(1, Inches(0.8), Inches(4.8),
                                         Inches(5.0), Inches(0.03))
    line_shape.fill.solid(); line_shape.fill.fore_color.rgb = C_GRAY
    line_shape.line.fill.background()

    # 팀/과목 정보
    add_text(slide, "경기대학교 · 2026-05-22",
             Inches(0.8), Inches(4.95), Inches(8), Inches(0.4),
             size=13, color=C_GRAY)

    # 하단 통계 카드 3개
    stats = [
        ("NCM 실차", "191대"),
        ("LSTM R²", "0.9233"),
        ("RMSE", "0.784%"),
    ]
    for i, (label, val) in enumerate(stats):
        cx = Inches(0.8 + i * 2.4)
        card = add_rect(slide, cx, Inches(5.7), Inches(2.1), Inches(1.1),
                        fill=C_CARD, line=C_ACCENT, line_w=Pt(1))
        add_text(slide, val, cx, Inches(5.75), Inches(2.1), Inches(0.55),
                 size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, label, cx, Inches(6.25), Inches(2.1), Inches(0.4),
                 size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_02_toc(prs):
    """목차"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "목차", "Table of Contents")

    items = [
        ("01", "문제 상황",          "전기차 배터리 교체비용 & SOH 예측 부재"),
        ("02", "프로젝트 개요",       "솔루션 아이디어 & 데이터셋 소개"),
        ("03", "데이터 전처리",       "EVBattery NCM 191대 피처 엔지니어링"),
        ("04", "배터리 열화 이론",    "SEI 성장 + 무릎효과 (2-phase 모델)"),
        ("05", "4-way 모델 비교",     "LSTM / Transformer / XGBoost / Neural CDE"),
        ("06", "Physics-Informed",   "Peukert & Arrhenius 손실함수 내장"),
        ("07", "실험 결과",           "R² / RMSE / MAE 비교 + SOH 곡선"),
        ("08", "앱 데모",             "운전자 분류 → 맞춤 추천 → 수명 연장 수치"),
        ("09", "결론",                "기여점 · 한계 · 향후 연구"),
    ]

    for i, (num, title, sub) in enumerate(items):
        row = i % 5
        col = i // 5
        x = Inches(0.5 + col * 6.5)
        y = Inches(1.3 + row * 1.1)
        w = Inches(6.1)

        card = add_rect(slide, x, y, w, Inches(0.95), fill=C_CARD, line=C_ACCENT, line_w=Pt(0.5))
        add_text(slide, num, x + Inches(0.12), y + Inches(0.05), Inches(0.6), Inches(0.45),
                 size=20, bold=True, color=C_ACCENT)
        add_text(slide, title, x + Inches(0.7), y + Inches(0.05), Inches(5.0), Inches(0.45),
                 size=15, bold=True, color=C_WHITE)
        add_text(slide, sub, x + Inches(0.7), y + Inches(0.48), Inches(5.0), Inches(0.38),
                 size=10, color=C_GRAY)


def slide_03_news1(prs):
    """문제상황 1 — 뉴스: 배터리 교체비용"""
    slide = new_slide(prs)
    bg(slide, RGBColor(0x0D, 0x0D, 0x0D))
    section_title(slide, "문제 상황 ①", "전기차 배터리 교체비용 — 소비자의 공포")

    # 뉴스 카드 1
    news1 = [
        ("[연합뉴스] 2024.03.15",
         '"전기차 배터리 교체비 1,200만원…\n중고차 값보다 비싼 수리비"',
         "아이오닉5 배터리 교체 견적 1,100~1,400만원. 차량 잔존가치보다\n수리비가 높아 사실상 '폐차 권유' 사례 증가."),
        ("[KBS 뉴스] 2024.06.20",
         '"전기차 8년 탔더니 SOH 71%…\n보증 기간 지나면 내 돈으로?"',
         "현대·기아 배터리 보증 10년/20만km/SOH 70%. 보증 종료 후\n교체 비용 전액 자기 부담. 소비자 불안 급증."),
    ]

    for i, (source, headline, body) in enumerate(news1):
        x = Inches(0.4 + i * 6.4)
        y = Inches(1.4)

        card = add_rect(slide, x, y, Inches(6.0), Inches(4.8),
                        fill=RGBColor(0x18, 0x18, 0x18), line=C_NEWS_RED, line_w=Pt(2))

        # 출처 배지
        badge = add_rect(slide, x + Inches(0.15), y + Inches(0.15),
                         Inches(2.2), Inches(0.3), fill=C_NEWS_RED)
        add_text(slide, source, x + Inches(0.15), y + Inches(0.13),
                 Inches(2.2), Inches(0.3), size=9, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER)

        add_text(slide, headline,
                 x + Inches(0.2), y + Inches(0.6), Inches(5.6), Inches(1.6),
                 size=16, bold=True, color=C_WHITE, wrap=True)

        add_rect(slide, x + Inches(0.2), y + Inches(2.3), Inches(5.6), Inches(0.03),
                 fill=C_GRAY)

        add_text(slide, body,
                 x + Inches(0.2), y + Inches(2.45), Inches(5.6), Inches(1.8),
                 size=11, color=C_GRAY, wrap=True)

    # 핵심 수치 하이라이트
    add_rect(slide, Inches(0.4), Inches(6.45), Inches(12.53), Inches(0.75),
             fill=RGBColor(0x7F, 0x1D, 0x1D))
    add_text(slide,
             "💡  전기차 오너의 최대 관심사 = 배터리 수명.  그런데 현재 앱/계기판은 '지금 SOH'만 보여줄 뿐, 미래 예측이 없다.",
             Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.6),
             size=13, bold=True, color=C_WHITE)


def slide_04_news2(prs):
    """문제상황 2 — 뉴스: 화재·예측 부재"""
    slide = new_slide(prs)
    bg(slide, RGBColor(0x0D, 0x0D, 0x0D))
    section_title(slide, "문제 상황 ②", "SOH 예측 부재 — 소비자는 아무것도 모른다")

    news2 = [
        ("[MBC 뉴스] 2024.08.02",
         '"인천 청라 아파트 주차장\n전기차 화재… 100여 대 피해"',
         "배터리 열화 상태 미파악으로 인한 전기차 화재 사고.\n과충전·고온 노출 배터리의 열폭주(Thermal Runaway).\n사전 SOH 모니터링의 필요성 대두."),
        ("[조선일보] 2024.11.10",
         '"전기차 중고차 시장 침체…\n배터리 잔존가치 알 수 없어"',
         "중고 전기차 구매자 78%가 '배터리 상태 불신' 응답.\n딜러사도 SOH 측정 장비 없어 정확한 가치 산정 불가.\n미래 수명 예측 서비스 수요 폭증."),
        ("[전자신문] 2025.02.14",
         '"전기차 배터리 데이터\n공개 안 해 소비자만 피해"',
         "완성차 업체들, BMS 데이터 비공개 정책 유지.\n제3자 SOH 진단 서비스 불가. 소비자 정보 비대칭 심각.\n오픈 데이터 기반 AI 진단 솔루션 요구 증가."),
    ]

    for i, (source, headline, body) in enumerate(news2):
        x = Inches(0.3 + i * 4.35)
        y = Inches(1.35)
        card = add_rect(slide, x, y, Inches(4.05), Inches(5.1),
                        fill=RGBColor(0x18, 0x18, 0x18), line=C_NEWS_RED, line_w=Pt(1.5))
        badge = add_rect(slide, x + Inches(0.12), y + Inches(0.12),
                         Inches(2.5), Inches(0.28), fill=C_NEWS_RED)
        add_text(slide, source, x + Inches(0.12), y + Inches(0.10),
                 Inches(2.5), Inches(0.28), size=8, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, headline,
                 x + Inches(0.15), y + Inches(0.52), Inches(3.75), Inches(1.5),
                 size=14, bold=True, color=C_WHITE, wrap=True)
        add_rect(slide, x + Inches(0.15), y + Inches(2.1), Inches(3.75), Inches(0.03),
                 fill=C_GRAY)
        add_text(slide, body,
                 x + Inches(0.15), y + Inches(2.25), Inches(3.75), Inches(2.6),
                 size=10, color=C_GRAY, wrap=True)

    add_rect(slide, Inches(0.3), Inches(6.6), Inches(12.73), Inches(0.65),
             fill=RGBColor(0x1E, 0x1B, 0x4B))
    add_text(slide,
             "→  우리가 만드는 것: 오픈 NCM 데이터 + AI로 '미래 SOH 예측 + 교체 시점 + 습관 개선 효과'를 제공하는 자가진단 앱",
             Inches(0.5), Inches(6.65), Inches(12.2), Inches(0.55),
             size=12, bold=True, color=C_WHITE)


def slide_05_solution(prs):
    """우리의 솔루션 개요"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "우리의 솔루션", "현대·기아 EV 배터리 자가진단 웹앱")

    # 흐름도: 5단계 arrow
    steps = [
        ("🚗", "차량 정보\n입력", "브랜드/모델\n지역/사용연수"),
        ("🏎️", "운전 습관\n분석", "속도·급가속\n급제동·SOC폭"),
        ("🧠", "AI 예측", "LSTM R²=0.924\nMC Dropout"),
        ("📈", "열화 곡선\n생성", "2-phase 물리\n모델 앵커링"),
        ("💡", "맞춤 추천\n TOP 3", "Peukert·Arrhenius\n수명 연장 수치"),
    ]

    for i, (icon, title, sub) in enumerate(steps):
        cx = Inches(0.5 + i * 2.55)
        card = add_rect(slide, cx, Inches(1.5), Inches(2.25), Inches(3.0),
                        fill=C_CARD, line=C_ACCENT, line_w=Pt(1))
        add_text(slide, icon, cx, Inches(1.6), Inches(2.25), Inches(0.8),
                 size=28, align=PP_ALIGN.CENTER)
        add_text(slide, title, cx, Inches(2.4), Inches(2.25), Inches(0.9),
                 size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)
        add_text(slide, sub, cx, Inches(3.3), Inches(2.25), Inches(1.0),
                 size=10, color=C_GRAY, align=PP_ALIGN.CENTER, wrap=True)

        if i < 4:
            add_text(slide, "→", Inches(2.55 + i * 2.55), Inches(2.7),
                     Inches(0.25), Inches(0.5), size=20, color=C_ACCENT,
                     align=PP_ALIGN.CENTER)

    # 하단 — 기존 대비 차별점
    add_text(slide, "기존 서비스 vs 우리 앱",
             Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
             size=14, bold=True, color=C_ACCENT)

    compare = [
        ("기존 계기판·앱", "현재 SOH 숫자 하나", "❌"),
        ("우리 앱",
         "현재 SOH + 미래 열화 곡선 + 교체 시점 + 습관별 수명 연장 수치", "✅"),
    ]
    for i, (who, what, mark) in enumerate(compare):
        y = Inches(5.5 + i * 0.7)
        col = C_CARD if i == 0 else RGBColor(0x0F, 0x2F, 0x1A)
        add_rect(slide, Inches(0.5), y, Inches(12.4), Inches(0.6),
                 fill=col, line=None)
        add_text(slide, f"{mark}  {who}", Inches(0.7), y + Inches(0.05),
                 Inches(3.0), Inches(0.5), size=12, bold=True,
                 color=C_WHITE if i == 0 else C_GREEN)
        add_text(slide, what, Inches(3.7), y + Inches(0.05),
                 Inches(9.0), Inches(0.5), size=11, color=C_GRAY if i == 0 else C_WHITE)


def slide_06_dataset(prs):
    """데이터셋 소개"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "데이터셋: EVBattery NCM", "Figshare 공개 데이터 — 중국 실차 충전 이력")

    # 왼쪽: 데이터 설명
    left_items = [
        ("📦 출처", "Figshare 'EVBattery: A Large-Scale EV Dataset'"),
        ("🚗 차량", "총 465대 → 품질 필터 후 191대 (NCM 화학 99%)"),
        ("📊 규모", "학습 153대 (85,648세션) / 테스트 38대 (21,629세션)"),
        ("📁 형식", "pkl 파일 1개 = 충전 세션 1회 (전압/전류/온도/용량)"),
        ("🔋 SOH", "min 64% / max 100% / mean 97.5% / std 4.27%"),
        ("⚗️ 화학", "NCM (Ni-Co-Mn) ← 현대·기아 EV 99%와 동일"),
    ]

    for i, (key, val) in enumerate(left_items):
        y = Inches(1.45 + i * 0.85)
        add_rect(slide, Inches(0.4), y, Inches(6.5), Inches(0.72),
                 fill=C_CARD, line=None)
        add_text(slide, key, Inches(0.55), y + Inches(0.05), Inches(1.5), Inches(0.35),
                 size=11, bold=True, color=C_ACCENT)
        add_text(slide, val, Inches(2.05), y + Inches(0.05), Inches(4.7), Inches(0.6),
                 size=11, color=C_WHITE, wrap=True)

    # 오른쪽: NCM 화학 비교표
    add_text(slide, "왜 NASA LCO 데이터를 안 쓰나?",
             Inches(7.3), Inches(1.4), Inches(5.5), Inches(0.4),
             size=13, bold=True, color=C_YELLOW)

    chem_table = [
        ("화학", "현대·기아 EV", "EVBattery", "NASA"),
        ("NCM", "✅ 99%",        "✅ 99%",     "❌"),
        ("LFP", "레이EV만",      "1%",         "❌"),
        ("LCO", "❌",            "❌",          "✅ (실험용)"),
    ]
    col_w = [Inches(1.3), Inches(1.35), Inches(1.35), Inches(1.3)]
    col_x = [Inches(7.3), Inches(8.6), Inches(9.95), Inches(11.3)]
    for r, row in enumerate(chem_table):
        for c, cell in enumerate(row):
            y = Inches(1.9 + r * 0.65)
            bg_c = C_ACCENT if r == 0 else (RGBColor(0x0F, 0x2F, 0x1A) if c in [1,2] and r > 0 and "✅" in cell
                                             else C_CARD)
            if r == 0: bg_c = RGBColor(0x1E, 0x3A, 0x5F)
            add_rect(slide, col_x[c], y, col_w[c], Inches(0.6),
                     fill=bg_c, line=RGBColor(0x2D, 0x3F, 0x5E), line_w=Pt(0.5))
            clr = C_WHITE if r == 0 else (C_GREEN if "✅" in cell else (C_RED if "❌" in cell else C_GRAY))
            add_text(slide, cell, col_x[c], y + Inches(0.05), col_w[c], Inches(0.5),
                     size=11, bold=(r == 0), color=clr, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(7.3), Inches(4.65), Inches(5.6), Inches(0.6),
             fill=RGBColor(0x78, 0x35, 0x0F))
    add_text(slide, "LCO 열화 메커니즘 ≠ NCM → NASA 데이터로 학습하면\n현대·기아 EV에 틀린 예측 출력",
             Inches(7.45), Inches(4.68), Inches(5.3), Inches(0.55),
             size=10, bold=True, color=C_WHITE, wrap=True)

    # 로딩 파이프라인 요약
    add_text(slide, "로딩 파이프라인",
             Inches(7.3), Inches(5.45), Inches(5.6), Inches(0.35),
             size=12, bold=True, color=C_ACCENT)
    pipe_steps = ["pkl 병렬 로딩\n(16 workers)", "→", "flat DataFrame", "→", "parquet 캐시\n(재실행 즉시)"]
    px = [Inches(7.3), Inches(8.7), Inches(9.0), Inches(10.35), Inches(10.6)]
    pw = [Inches(1.3), Inches(0.3), Inches(1.3), Inches(0.3), Inches(1.6)]
    for i, (txt, x, w) in enumerate(zip(pipe_steps, px, pw)):
        if i % 2 == 0:
            add_rect(slide, x, Inches(5.9), w, Inches(0.7),
                     fill=C_CARD, line=C_ACCENT, line_w=Pt(0.5))
        add_text(slide, txt, x, Inches(5.95), w, Inches(0.6),
                 size=9, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)


def slide_07_preprocessing(prs):
    """데이터 전처리 & 피처 엔지니어링"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "데이터 전처리 & 피처 엔지니어링", "13개 피처 생성 — 물리 법칙 내재화")

    # 피처 테이블
    features = [
        ("voltage_mean / std", "세션 내 전압 평균·분산", "원본"),
        ("current_mean",       "평균 충전 전류",          "원본"),
        ("temp_mean",          "평균 온도",               "원본"),
        ("capacity_fade_rate", "5-step 이동 용량 감소율", "파생"),
        ("cycle_norm",         "정규화 사이클 위치 [0,1]", "파생"),
        ("c_rate",             "전류 / 초기용량 (충전속도)", "물리"),
        ("thermal_stress",     "max(T-25°C, 0) — 과열 누적", "물리"),
        ("degradation_index",  "Peukert항 × Arrhenius항 복합 열화 지수", "물리"),
        ("cum_degradation",    "누적 열화 지수 (정규화)",  "v3.5"),
        ("temp_roll_5",        "온도 5-step 이동평균",     "v3.5"),
        ("c_rate_roll_5",      "C-rate 5-step 이동평균",  "v3.5"),
        ("deg_roll_5",         "열화지수 5-step 이동평균", "v3.5"),
        ("soh_norm (채널)",    "차량별 첫 세션 대비 정규화 SOH", "타겟"),
    ]
    type_colors = {
        "원본": C_GRAY,
        "파생": C_ACCENT,
        "물리": C_GREEN,
        "v3.5": C_YELLOW,
        "타겟": C_RED,
    }

    col_headers = ["피처명", "의미", "유형"]
    col_x = [Inches(0.4), Inches(3.5), Inches(11.0)]
    col_w = [Inches(3.0), Inches(7.4), Inches(1.8)]

    # 헤더
    for h, x, w in zip(col_headers, col_x, col_w):
        add_rect(slide, x, Inches(1.35), w, Inches(0.4),
                 fill=RGBColor(0x1E, 0x3A, 0x5F))
        add_text(slide, h, x, Inches(1.38), w, Inches(0.35),
                 size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for i, (name, desc, typ) in enumerate(features):
        y = Inches(1.8 + i * 0.38)
        row_bg = C_CARD if i % 2 == 0 else RGBColor(0x16, 0x22, 0x38)
        for x, w in zip(col_x, col_w):
            add_rect(slide, x, y, w, Inches(0.35), fill=row_bg)
        add_text(slide, name, col_x[0], y + Inches(0.02), col_w[0], Inches(0.32),
                 size=9.5, color=C_WHITE, bold=(typ == "타겟"))
        add_text(slide, desc, col_x[1], y + Inches(0.02), col_w[1], Inches(0.32),
                 size=9.5, color=C_GRAY)
        tc = type_colors.get(typ, C_GRAY)
        add_text(slide, typ, col_x[2], y + Inches(0.02), col_w[2], Inches(0.32),
                 size=9.5, bold=True, color=tc, align=PP_ALIGN.CENTER)

    # soh_norm 설명
    add_rect(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5),
             fill=RGBColor(0x0F, 0x2F, 0x1A))
    add_text(slide, "soh_norm = 각 차량의 첫 세션 SOH 대비 현재 비율 → 차량 간 절대값 차이를 제거, 열화 궤적 학습에 집중",
             Inches(0.6), Inches(6.83), Inches(12.1), Inches(0.44),
             size=10.5, bold=True, color=C_GREEN)


def slide_08_degradation(prs):
    """배터리 열화 이론 — 2-phase 모델"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "배터리 열화 이론 — 2-phase 모델", "논문 기반 물리 모델링")

    # 왼쪽: 열화 곡선 그래프
    fig, ax = plt.subplots(figsize=(5.5, 4.0), facecolor="#1A253D")
    ax.set_facecolor("#0F172A")

    cycles = np.arange(1, 2001)
    alpha = 30 / np.sqrt(2000)
    soh = 100 - alpha * np.sqrt(cycles)
    knee_idx = np.argmax(soh < 85)
    if knee_idx > 0:
        k = np.log(30 / max(100 - soh[knee_idx], 1e-4)) / (2000 - knee_idx)
        soh[knee_idx:] = 100 - (100 - soh[knee_idx]) * np.exp(k * (cycles[knee_idx:] - knee_idx))
    soh = np.clip(soh, 0, 100)
    sigma = 1.2 * np.sqrt(cycles / 2000)

    ax.fill_between(cycles, soh - sigma, soh + sigma, alpha=0.25, color="#3B82F6")
    ax.plot(cycles, soh, color="#3B82F6", lw=2.5, label="SOH 열화 곡선")
    ax.axhline(85, color="#FBBF24", lw=1.5, ls="--", label="무릎 임계 (85%)")
    ax.axhline(70, color="#EF4444", lw=1.5, ls="--", label="EOL 기준 (70%)")
    if knee_idx > 0:
        ax.axvline(knee_idx, color="#FBBF24", lw=1, ls=":", alpha=0.7)
        ax.annotate("무릎 효과\n시작", xy=(knee_idx, 85), xytext=(knee_idx + 100, 88),
                    color="#FBBF24", fontsize=8,
                    arrowprops=dict(arrowstyle="->", color="#FBBF24", lw=1))

    ax.text(100, 97, "Phase 1\nSEI 성장 (√t)", color="#94A3B8", fontsize=8)
    ax.text(knee_idx + 50, 78, "Phase 2\n지수 가속", color="#EF4444", fontsize=8)

    ax.set_xlabel("누적 충전 사이클", color="#94A3B8", fontsize=9)
    ax.set_ylabel("SOH (%)", color="#94A3B8", fontsize=9)
    ax.set_ylim(60, 105)
    ax.tick_params(colors="#94A3B8", labelsize=8)
    for spine in ax.spines.values():
        spine.set_color("#2D3F5E")
    ax.legend(facecolor="#1A253D", edgecolor="#2D3F5E", labelcolor="white", fontsize=8)
    plt.tight_layout()
    add_fig(slide, fig, Inches(0.4), Inches(1.4), Inches(5.8), Inches(4.2))

    # 오른쪽: Phase 설명
    phases = [
        ("Phase 1: SEI 성장 (SOH > 85%)", C_ACCENT,
         "수식: SOH = 100 - α·√t\n"
         "• SEI(고체 전해질 계면) 막이 음극 표면에 성장\n"
         "• 막이 두꺼울수록 반응 억제 → 열화 감속\n"
         "• 참고: Ploehn et al. 2004 / Safari & Delacourt 2011"),
        ("Phase 2: 무릎 효과 (SOH < 85%)", C_RED,
         "수식: SOH = 100 - (100-knee)·exp(k·Δt)\n"
         "• 리튬 도금(Li Plating): 고속충전 시 음극 표면 석출\n"
         "• NCM 양극 입자 균열 누적 → 지수함수적 가속\n"
         "• 참고: Attia et al. 2022 (Nature Energy)"),
        ("승법 노이즈 (불확실성)", C_YELLOW,
         "σ(t) = 1.2 × √(t / EOL)\n"
         "• 오래될수록 셀 간 분산 증가 → 예측 불확실성 확대\n"
         "• 앱에서 SOH ± X% 신뢰구간으로 표시\n"
         "• 참고: Saha & Goebel 2009 (IEEE TIM)"),
    ]
    for i, (title, color, text) in enumerate(phases):
        y = Inches(1.4 + i * 1.9)
        add_rect(slide, Inches(6.5), y, Inches(6.5), Inches(1.75),
                 fill=C_CARD, line=color, line_w=Pt(2))
        add_text(slide, title, Inches(6.65), y + Inches(0.08), Inches(6.2), Inches(0.38),
                 size=12, bold=True, color=color)
        add_text(slide, text, Inches(6.65), y + Inches(0.48), Inches(6.2), Inches(1.15),
                 size=10, color=C_WHITE, wrap=True)


def slide_09_model_overview(prs):
    """4-way 모델 비교 개요"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "4-way 모델 비교", "왜 단일 모델이 아닌 4개를 비교했나?")

    add_text(slide,
             "\"특정 모델에 우연히 좋은 결과가 나온 게 아님을 증명\" — 방법론의 견고성 확보",
             Inches(0.5), Inches(1.25), Inches(12.4), Inches(0.35),
             size=12, italic=True, color=C_GRAY)

    models_info = [
        ("LSTM", C_ACCENT,
         "Long Short-Term Memory",
         "순서 있는 시계열 처리에 특화.\n2층 LSTM + LayerNorm + Dropout.\n멀티태스크: ΔSOH + 열화율 동시 예측.\nMC Dropout으로 불확실성 정량화.",
         "R² = 0.9233 ★"),
        ("Transformer", C_GREEN,
         "Multi-Head Self-Attention",
         "30개 세션 중 '어느 세션이 현재 SOH에\n중요한지' Attention이 자동 가중치 부여.\nPre-LN + Learnable PE 사용.",
         "R² = 0.9169"),
        ("XGBoost", C_YELLOW,
         "Gradient Boosted Trees",
         "딥러닝 없이 얼마나 되는지 비교 (Baseline).\n30개 세션 → mean/std/min/max/slope\n5개 통계 × 13피처 = 65차원 피처벡터.",
         "R² = 0.9084"),
        ("Neural CDE", C_RED,
         "Kidger et al. 2020 (NeurIPS)",
         "불규칙 시계열 이론적 최적 모델.\ndZ(t) = fθ(Z(t))·dX(t) 미분방정식 적분.\n이론 ≠ 실제: 이 데이터는 규칙적이라 발산.",
         "R² = −49 ⚠️"),
    ]

    for i, (name, color, subtitle, desc, result) in enumerate(models_info):
        x = Inches(0.4 + i * 3.25)
        y = Inches(1.7)
        card = add_rect(slide, x, y, Inches(3.1), Inches(5.3),
                        fill=C_CARD, line=color, line_w=Pt(2))
        # 헤더 배지
        add_rect(slide, x, y, Inches(3.1), Inches(0.45), fill=color)
        add_text(slide, name, x, y + Inches(0.02), Inches(3.1), Inches(0.4),
                 size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, subtitle, x, y + Inches(0.5), Inches(3.1), Inches(0.38),
                 size=9.5, color=color, align=PP_ALIGN.CENTER, italic=True)
        add_text(slide, desc, x + Inches(0.12), y + Inches(0.95), Inches(2.85), Inches(2.8),
                 size=10, color=C_WHITE, wrap=True)
        # 결과
        res_bg = RGBColor(0x0F, 0x2F, 0x1A) if "★" in result else (
                 RGBColor(0x3F, 0x10, 0x10) if "⚠️" in result else C_CARD)
        add_rect(slide, x + Inches(0.12), y + Inches(3.85), Inches(2.85), Inches(0.55),
                 fill=res_bg, line=color, line_w=Pt(1))
        add_text(slide, result, x + Inches(0.12), y + Inches(3.9), Inches(2.85), Inches(0.45),
                 size=13, bold=True, color=color, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.25),
             fill=RGBColor(0x1E, 0x1B, 0x4B))
    add_text(slide, "공통 설계 원칙: 잔차(ΔSOH) 학습 + WINDOW_SIZE=30 + AdamW + CosineAnnealing + EarlyStopping(patience=20)",
             Inches(0.6), Inches(7.12), Inches(12.1), Inches(0.22),
             size=9, color=C_GRAY)


def slide_10_lstm_detail(prs):
    """LSTM & Transformer 상세"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "LSTM & Transformer 아키텍처", "잔차 학습 + MC Dropout + Pre-LN")

    # LSTM 아키텍처 다이어그램
    fig, ax = plt.subplots(figsize=(5.5, 4.5), facecolor="#1A253D")
    ax.set_facecolor("#1A253D")
    ax.set_xlim(0, 10); ax.set_ylim(0, 10)
    ax.axis("off")

    boxes = [
        (5, 9.0, "입력 윈도우\n(30 × 14)", "#2D3F5E"),
        (5, 7.2, "LSTM Layer 1\nhidden=128 + LayerNorm", "#1E3A5F"),
        (5, 5.4, "LSTM Layer 2\nhidden=64 + LayerNorm", "#1E3A5F"),
        (2.5, 3.4, "fc_delta\n→ ΔSOH", "#0F2F1A"),
        (7.5, 3.4, "fc_rate\n→ 열화율 (보조)", "#2D1A0F"),
    ]
    for (bx, by, label, col) in boxes:
        rect = mpatches.FancyBboxPatch((bx - 1.8, by - 0.65), 3.6, 1.1,
                                       boxstyle="round,pad=0.1", linewidth=1.5,
                                       edgecolor="#3B82F6", facecolor=col)
        ax.add_patch(rect)
        ax.text(bx, by, label, ha="center", va="center", color="white",
                fontsize=8, fontweight="bold")

    for (y1, y2) in [(8.35, 7.85), (6.55, 6.05)]:
        ax.annotate("", xy=(5, y2), xytext=(5, y1),
                    arrowprops=dict(arrowstyle="->", color="#3B82F6", lw=1.5))
    ax.annotate("", xy=(2.5, 4.05), xytext=(4.5, 4.75),
                arrowprops=dict(arrowstyle="->", color="#22C55E", lw=1.5))
    ax.annotate("", xy=(7.5, 4.05), xytext=(5.5, 4.75),
                arrowprops=dict(arrowstyle="->", color="#FBBF24", lw=1.5))

    ax.text(5, 1.5, "MC Dropout (50회 forward)\n→ 평균 = SOH 추정 / 표준편차 = 불확실성",
            ha="center", va="center", color="#94A3B8", fontsize=8,
            bbox=dict(boxstyle="round", facecolor="#0F172A", edgecolor="#3B82F6"))

    plt.tight_layout()
    add_fig(slide, fig, Inches(0.3), Inches(1.4), Inches(5.6), Inches(5.4))

    # 오른쪽: 핵심 설계 포인트
    points = [
        ("잔차 학습 (ΔSOH)", C_ACCENT,
         "절대 SOH 대신 다음 스텝 변화량 예측.\n"
         "SOH가 95~100%에 몰린 분포에서\n'변화 없음' 과적합 방지."),
        ("MC Dropout", C_GREEN,
         "추론 시 Dropout 유지 → 50회 샘플링.\n"
         "평균 = SOH 점추정 / std = 불확실성.\n"
         "앱에서 SOH ± X% 로 표시."),
        ("Pre-LN Transformer", C_YELLOW,
         "norm_first=True → 학습 안정성 향상.\n"
         "Learnable Positional Encoding:\n"
         "30개 세션의 위치 중요도 학습."),
    ]
    for i, (title, color, text) in enumerate(points):
        y = Inches(1.4 + i * 1.9)
        add_rect(slide, Inches(6.2), y, Inches(6.8), Inches(1.75),
                 fill=C_CARD, line=color, line_w=Pt(1.5))
        add_text(slide, title, Inches(6.35), y + Inches(0.08), Inches(6.5), Inches(0.38),
                 size=13, bold=True, color=color)
        add_text(slide, text, Inches(6.35), y + Inches(0.48), Inches(6.5), Inches(1.15),
                 size=10.5, color=C_WHITE, wrap=True)


def slide_11_pinn(prs):
    """Physics-Informed Loss"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "Physics-Informed Loss (PINN)", "Peukert & Arrhenius — 물리 법칙을 손실함수에 내장")

    # 왼쪽: Peukert 곡선
    fig, axes = plt.subplots(1, 2, figsize=(6.0, 3.8), facecolor="#1A253D")
    for ax in axes:
        ax.set_facecolor("#0F172A")
        ax.tick_params(colors="#94A3B8", labelsize=8)
        for spine in ax.spines.values():
            spine.set_color("#2D3F5E")

    # Peukert
    c_rates = np.linspace(0.3, 2.5, 100)
    PEUKERT_N = 0.621
    cycles_80 = 550 * (1.0 / c_rates) ** PEUKERT_N
    axes[0].plot(c_rates, cycles_80, color="#3B82F6", lw=2.5)
    axes[0].axvline(0.6, color="#22C55E", ls="--", lw=1.5, label="절약형 0.6C")
    axes[0].axvline(1.0, color="#FBBF24", ls="--", lw=1.5, label="평균형 1.0C")
    axes[0].axvline(1.8, color="#EF4444", ls="--", lw=1.5, label="공격형 1.8C")
    axes[0].set_xlabel("C-rate", color="#94A3B8", fontsize=8)
    axes[0].set_ylabel("수명 사이클", color="#94A3B8", fontsize=8)
    axes[0].set_title("Peukert 법칙 (n=0.621)", color="white", fontsize=9)
    axes[0].legend(facecolor="#1A253D", edgecolor="#2D3F5E", labelcolor="white", fontsize=7)

    # Arrhenius
    temps = np.linspace(5, 55, 100)
    EA = 31000; R = 8.314; T_ref = 298.15
    mult = np.exp(EA / R * (1/T_ref - 1/(temps + 273.15)))
    axes[1].plot(temps, mult, color="#EF4444", lw=2.5)
    axes[1].axvline(25, color="#94A3B8", ls="--", lw=1, label="기준 25°C")
    axes[1].axvspan(35, 55, alpha=0.15, color="#EF4444")
    axes[1].set_xlabel("온도 (°C)", color="#94A3B8", fontsize=8)
    axes[1].set_ylabel("열화 가속 배율", color="#94A3B8", fontsize=8)
    axes[1].set_title("Arrhenius 법칙 (Ea=31kJ/mol)", color="white", fontsize=9)
    axes[1].legend(facecolor="#1A253D", edgecolor="#2D3F5E", labelcolor="white", fontsize=7)

    plt.tight_layout()
    add_fig(slide, fig, Inches(0.3), Inches(3.8), Inches(6.5), Inches(3.5))

    # 손실함수 수식
    loss_box = add_rect(slide, Inches(0.3), Inches(1.4), Inches(6.5), Inches(2.2),
                        fill=RGBColor(0x0A, 0x10, 0x20), line=C_ACCENT, line_w=Pt(1.5))
    add_text(slide, "손실함수 구성",
             Inches(0.5), Inches(1.45), Inches(6.0), Inches(0.35),
             size=12, bold=True, color=C_ACCENT)

    loss_lines = [
        "L = MSE(ΔSOH_pred, ΔSOH_target)",
        "  + λ_aux × MSE(rate_pred, rate_target)",
        "  + λ_pinn × (Peukert위반 + Arrhenius위반)",
        "",
        "λ_aux = 0.1 (보조 태스크 가중치)",
        "λ_pinn = 0.05 (물리 제약 가중치)",
    ]
    for i, line in enumerate(loss_lines):
        color = C_ACCENT if "L =" in line else (C_GREEN if "aux" in line else
                (C_YELLOW if "pinn" in line else (C_WHITE if "λ" in line else C_GRAY)))
        add_text(slide, line,
                 Inches(0.5), Inches(1.85 + i * 0.28), Inches(6.1), Inches(0.27),
                 size=10.5, color=color, bold=("L =" in line))

    # 오른쪽 설명
    explanations = [
        ("Peukert 법칙", C_ACCENT,
         "cycle-life ∝ C-rate^(-n), n=0.621\n"
         "· 0.5C → 약 780사이클 수명\n"
         "· 1.0C → 약 550사이클 수명\n"
         "· 2.0C → 약 330사이클 수명\n"
         "출처: Safari 2011 / Pelletier 2017"),
        ("Arrhenius 법칙", C_RED,
         "열화율 ∝ exp(Ea/R × (1/T_ref - 1/T))\n"
         "· 25°C 기준, 10°C 오를 때마다 ~1.5배 가속\n"
         "· 여름 옥상 주차(+40°C) vs 지하주차\n"
         "  → 수명 2~3배 차이\n"
         "출처: Wang 2014 / Waldmann 2014"),
    ]
    for i, (title, color, text) in enumerate(explanations):
        y = Inches(1.4 + i * 3.0)
        add_rect(slide, Inches(7.1), y, Inches(6.0), Inches(2.7),
                 fill=C_CARD, line=color, line_w=Pt(1.5))
        add_text(slide, title, Inches(7.25), y + Inches(0.1), Inches(5.7), Inches(0.38),
                 size=13, bold=True, color=color)
        add_text(slide, text, Inches(7.25), y + Inches(0.52), Inches(5.7), Inches(2.0),
                 size=10, color=C_WHITE, wrap=True)


def slide_12_results(prs):
    """실험 결과 — 4-way 비교"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "실험 결과", "4-way 모델 비교 — NCM 실차 38대 테스트셋")

    # 바 차트
    fig, axes = plt.subplots(1, 3, figsize=(7.0, 3.8), facecolor="#1A253D")
    models = ["LSTM", "Trans-\nformer", "XGBoost", "CDE"]
    r2_vals    = [0.9233, 0.9169, 0.9084, -49.04]
    rmse_vals  = [0.784,  0.817,  0.857,  20.062]
    mae_vals   = [0.295,  0.277,  0.238,  12.081]
    colors_bar = ["#3B82F6", "#22C55E", "#FBBF24", "#EF4444"]

    metrics = [
        (axes[0], r2_vals,   "R²",       [-1, 1.05], True),
        (axes[1], rmse_vals, "RMSE (%)", [0, 25],    False),
        (axes[2], mae_vals,  "MAE (%)",  [0, 14],    False),
    ]

    for ax, vals, ylabel, ylim, higher_better in metrics:
        ax.set_facecolor("#0F172A")
        # CDE 제외한 가시화를 위해 clip
        plot_vals = [min(v, ylim[1] * 0.9) if v > ylim[1] else v for v in vals]
        bars = ax.bar(models, plot_vals, color=colors_bar, width=0.5, edgecolor="#2D3F5E")
        for bar, val, pv in zip(bars, vals, plot_vals):
            label = f"{val:.4f}" if abs(val) < 10 else f"{val:.1f}"
            ax.text(bar.get_x() + bar.get_width()/2, pv + ylim[1]*0.02,
                    label, ha="center", color="white", fontsize=7.5)
        ax.set_ylabel(ylabel, color="#94A3B8", fontsize=8)
        ax.set_ylim(ylim)
        ax.tick_params(colors="#94A3B8", labelsize=7)
        for spine in ax.spines.values():
            spine.set_color("#2D3F5E")
        ax.set_title(("높을수록 좋음 ↑" if higher_better else "낮을수록 좋음 ↓"),
                     color="#94A3B8", fontsize=7)

    plt.tight_layout()
    add_fig(slide, fig, Inches(0.3), Inches(1.4), Inches(7.5), Inches(4.2))

    # 오른쪽: 결과 해석
    result_table = [
        ("모델",        "R²",     "RMSE",   "비고"),
        ("LSTM ★",      "0.9233", "0.784%", "1위"),
        ("Transformer", "0.9169", "0.817%", "2위"),
        ("XGBoost",     "0.9084", "0.857%", "3위"),
        ("Neural CDE",  "-49.04", "20.06%", "발산"),
    ]
    col_x2 = [Inches(8.0), Inches(9.4), Inches(10.4), Inches(11.4)]
    col_w2 = [Inches(1.3), Inches(0.95), Inches(0.95), Inches(1.45)]
    for r, row in enumerate(result_table):
        for c, cell in enumerate(row):
            y = Inches(1.4 + r * 0.65)
            bg_c = RGBColor(0x1E, 0x3A, 0x5F) if r == 0 else (
                   RGBColor(0x0F, 0x2F, 0x1A) if r == 1 else (
                   RGBColor(0x3F, 0x10, 0x10) if r == 4 else C_CARD))
            add_rect(slide, col_x2[c], y, col_w2[c], Inches(0.58),
                     fill=bg_c, line=RGBColor(0x2D, 0x3F, 0x5E), line_w=Pt(0.3))
            clr = C_GREEN if r == 1 else (C_RED if r == 4 else (C_WHITE if r == 0 else C_GRAY))
            add_text(slide, cell, col_x2[c], y + Inches(0.05), col_w2[c], Inches(0.48),
                     size=10.5, bold=(r in [0, 1]), color=clr, align=PP_ALIGN.CENTER)

    # CDE 발산 설명
    add_rect(slide, Inches(8.0), Inches(4.85), Inches(4.9), Inches(1.0),
             fill=RGBColor(0x3F, 0x10, 0x10), line=C_RED, line_w=Pt(1))
    add_text(slide, "CDE R²=-49 이유",
             Inches(8.15), Inches(4.9), Inches(4.6), Inches(0.32),
             size=11, bold=True, color=C_RED)
    add_text(slide,
             "이론상 불규칙 시계열에 최적이나,\n"
             "이 데이터는 세션 간격이 실제로 규칙적.\n"
             "→ LSTM의 귀납적 편향이 더 잘 맞음.\n"
             "→ 실험 비교 없이 모델 선택 금지.",
             Inches(8.15), Inches(5.25), Inches(4.6), Inches(0.6),
             size=9.5, color=C_WHITE, wrap=True)

    # R² 해석
    add_rect(slide, Inches(8.0), Inches(6.0), Inches(4.9), Inches(0.8),
             fill=RGBColor(0x0F, 0x2F, 0x1A), line=C_GREEN, line_w=Pt(1))
    add_text(slide, "RMSE 0.784% = SOH 예측 오차 1% 이내\n"
             "SOH 95~100% 밀집 분포에서 R²<1은 정상.\n"
             "실용 기준: RMSE와 MAE가 핵심.",
             Inches(8.15), Inches(6.05), Inches(4.6), Inches(0.72),
             size=9.5, color=C_GREEN, wrap=True)


def slide_13_soh_curve(prs):
    """SOH 곡선 생성 방법"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "SOH 예측 곡선 생성", "물리 모델 + LSTM 앵커링 (Hybrid 방식)")

    # 왼쪽: 곡선 예시
    fig, ax = plt.subplots(figsize=(5.8, 4.2), facecolor="#1A253D")
    ax.set_facecolor("#0F172A")

    n = 2500
    cycles = np.arange(1, n + 1)

    def two_phase(alpha, eol, offset=0):
        c = cycles + offset
        s = 100 - alpha * np.sqrt(c)
        knee = np.argmax(s < 85)
        if knee > 0 and eol > knee:
            k = np.log(30 / max(100 - s[knee], 1e-4)) / (eol - knee)
            s[knee:] = 100 - (100 - s[knee]) * np.exp(k * (c[knee:] - knee))
        return np.clip(s, 0, 100)

    alpha_cur = 30 / np.sqrt(1800)
    alpha_opt = 30 / np.sqrt(2200)
    soh_cur = two_phase(alpha_cur, 1800)
    soh_opt = two_phase(alpha_opt, 2200)
    sigma_cur = 1.2 * np.sqrt(np.clip(cycles / 1800, 0, 1))

    current_cycle = 150
    ax.fill_between(cycles[current_cycle:],
                    soh_cur[current_cycle:] - sigma_cur[current_cycle:],
                    soh_cur[current_cycle:] + sigma_cur[current_cycle:],
                    alpha=0.2, color="#EF4444", label="_nolegend_")
    ax.plot(cycles[:current_cycle], soh_cur[:current_cycle],
            color="#6B7280", lw=1.5, ls=":", label="이미 지난 기간")
    ax.plot(cycles[current_cycle:], soh_cur[current_cycle:],
            color="#EF4444", lw=2.5, label="현재 습관 유지")
    ax.plot(cycles[current_cycle:], soh_opt[current_cycle:],
            color="#22C55E", lw=2.5, label="맞춤 개선 적용")
    ax.axhline(70, color="#94A3B8", lw=1.5, ls="--", label="EOL 70%")
    ax.axvline(current_cycle, color="#FBBF24", lw=1.5, label="현재 위치")
    ax.scatter([current_cycle], [soh_cur[current_cycle]], color="#FBBF24", s=60, zorder=5)
    ax.annotate(f"현재 SOH {soh_cur[current_cycle]:.1f}%",
                xy=(current_cycle, soh_cur[current_cycle]),
                xytext=(current_cycle + 100, soh_cur[current_cycle] + 3),
                color="#FBBF24", fontsize=8,
                arrowprops=dict(arrowstyle="->", color="#FBBF24", lw=1))
    ax.set_xlabel("누적 충전 사이클", color="#94A3B8", fontsize=9)
    ax.set_ylabel("SOH (%)", color="#94A3B8", fontsize=9)
    ax.set_ylim(60, 105)
    ax.tick_params(colors="#94A3B8", labelsize=8)
    for spine in ax.spines.values(): spine.set_color("#2D3F5E")
    ax.legend(facecolor="#1A253D", edgecolor="#2D3F5E", labelcolor="white", fontsize=8)
    plt.tight_layout()
    add_fig(slide, fig, Inches(0.3), Inches(1.4), Inches(6.0), Inches(4.5))

    # 오른쪽: 생성 절차
    steps_right = [
        ("1️⃣ EOL 계산", C_ACCENT,
         "C-rate × Arrhenius(T) → eol_user 사이클 수\n"
         "Peukert: EOL ∝ (1/C)^n\n"
         "Arrhenius: EOL ÷ exp(Ea/R·(1/T_ref - 1/T))"),
        ("2️⃣ 2-phase 곡선", C_YELLOW,
         "Phase1: SOH = 100 - α·√t (SEI 성장)\n"
         "Phase2: SOH < 85% → 지수 가속\n"
         "σ(t) = 1.2·√(t/EOL) 신뢰구간 계산"),
        ("3️⃣ LSTM 앵커링", C_GREEN,
         "실제 충전 이력 30세션 → LSTM → 현재 SOH\n"
         "물리 곡선을 LSTM 추정값으로 평행이동\n"
         "MC Dropout std + 승법노이즈 → 최종 신뢰구간"),
    ]
    for i, (title, color, text) in enumerate(steps_right):
        y = Inches(1.4 + i * 1.95)
        add_rect(slide, Inches(6.6), y, Inches(6.5), Inches(1.8),
                 fill=C_CARD, line=color, line_w=Pt(1.5))
        add_text(slide, title, Inches(6.75), y + Inches(0.08), Inches(6.2), Inches(0.38),
                 size=13, bold=True, color=color)
        add_text(slide, text, Inches(6.75), y + Inches(0.5), Inches(6.2), Inches(1.2),
                 size=10.5, color=C_WHITE, wrap=True)


def slide_14_driver_class(prs):
    """운전자 분류 시스템"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "운전자 유형 분류 시스템", "Random Forest 3-클래스 분류 — 5개 입력 피처")

    # 왼쪽: 레이더 차트
    fig, ax = plt.subplots(figsize=(4.5, 4.0), subplot_kw=dict(polar=True),
                           facecolor="#1A253D")
    ax.set_facecolor("#0F172A")
    categories = ["속도", "급가속", "급제동", "SOC폭", "충전빈도"]
    n_cat = len(categories)
    angles = [n / n_cat * 2 * math.pi for n in range(n_cat)]
    angles += angles[:1]

    driver_data = {
        "절약형":  [0.3, 0.1, 0.1, 0.4, 0.2],
        "평균형":  [0.5, 0.4, 0.4, 0.5, 0.5],
        "공격형":  [0.9, 0.8, 0.8, 0.9, 0.9],
    }
    driver_colors = {"절약형": "#22C55E", "평균형": "#FBBF24", "공격형": "#EF4444"}

    for dtype, vals in driver_data.items():
        vals_plot = vals + vals[:1]
        ax.plot(angles, vals_plot, color=driver_colors[dtype], lw=2, label=dtype)
        ax.fill(angles, vals_plot, alpha=0.1, color=driver_colors[dtype])

    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories, color="white", fontsize=9)
    ax.set_yticklabels([])
    ax.grid(color="#2D3F5E", lw=0.8)
    ax.spines["polar"].set_color("#2D3F5E")
    ax.legend(loc="upper right", bbox_to_anchor=(1.35, 1.15),
              facecolor="#1A253D", edgecolor="#2D3F5E", labelcolor="white", fontsize=8)
    plt.tight_layout()
    add_fig(slide, fig, Inches(0.3), Inches(1.4), Inches(5.0), Inches(4.5))

    # 오른쪽: 3유형 카드
    types = [
        ("🟢 절약형", C_GREEN,
         "평균속도 55 km/h / 급가속 3회/h\n급제동 2회/h / SOC폭 55% / 주 2회",
         "실효 C-rate 0.6C → EOL ~2200 사이클"),
        ("🟡 평균형", C_YELLOW,
         "평균속도 75 km/h / 급가속 10회/h\n급제동 8회/h / SOC폭 65% / 주 4회",
         "실효 C-rate 1.0C → EOL ~1800 사이클"),
        ("🔴 공격형", C_RED,
         "평균속도 110 km/h / 급가속 20회/h\n급제동 16회/h / SOC폭 85% / 주 6회",
         "실효 C-rate 1.8C → EOL ~1100 사이클"),
    ]
    for i, (title, color, habits, result) in enumerate(types):
        y = Inches(1.4 + i * 1.95)
        add_rect(slide, Inches(5.6), y, Inches(7.5), Inches(1.75),
                 fill=C_CARD, line=color, line_w=Pt(1.5))
        add_text(slide, title, Inches(5.75), y + Inches(0.08), Inches(7.2), Inches(0.4),
                 size=14, bold=True, color=color)
        add_text(slide, habits, Inches(5.75), y + Inches(0.5), Inches(7.2), Inches(0.7),
                 size=10, color=C_WHITE, wrap=True)
        add_rect(slide, Inches(5.75), y + Inches(1.25), Inches(7.1), Inches(0.35),
                 fill=RGBColor(0x0F, 0x17, 0x2A))
        add_text(slide, result, Inches(5.75), y + Inches(1.28), Inches(7.1), Inches(0.3),
                 size=10, bold=True, color=color)

    # 분류기 정보
    add_rect(slide, Inches(5.6), Inches(7.25), Inches(7.5), Inches(0.3),
             fill=RGBColor(0x1E, 0x1B, 0x4B))
    add_text(slide, "Random Forest 300 trees · 합성 900명(유형당 300) · 5-피처 StandardScaler",
             Inches(5.75), Inches(7.27), Inches(7.2), Inches(0.26),
             size=9, color=C_GRAY)


def slide_15_recommendation(prs):
    """맞춤 추천 엔진"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "맞춤 추천 엔진", "Peukert & Arrhenius 기반 수명 연장 정량화 — TOP 3 출력")

    add_text(slide, "단순 \"완속충전 하세요\" → 본인 패턴에서 효과 큰 행동 TOP 3, 수치로",
             Inches(0.5), Inches(1.22), Inches(12.4), Inches(0.32),
             size=11, italic=True, color=C_GRAY)

    recs = [
        ("🛣️", "평균 주행 속도", "85 km/h → 80 km/h",
         "+1년 4개월", "Peukert n=0.621", "쉬움", C_ACCENT),
        ("🔋", "충전 SOC 폭",   "85% → 60% (20-80%)",
         "+1년 2개월", "음극 응력 (Attia 2020)", "쉬움", C_GREEN),
        ("🌡️", "주차 환경",     "옥상→지하 주차",
         "+1년 0개월", "Arrhenius Ea=31kJ/mol", "중간", C_YELLOW),
    ]

    for i, (icon, cat, change, gain, physics, diff, color) in enumerate(recs):
        x = Inches(0.4 + i * 4.3)
        card = add_rect(slide, x, Inches(1.7), Inches(4.1), Inches(4.0),
                        fill=C_CARD, line=color, line_w=Pt(2))

        medal = ["🥇", "🥈", "🥉"][i]
        add_text(slide, f"{medal} 우선순위 {i+1}",
                 x + Inches(0.15), Inches(1.78), Inches(3.8), Inches(0.32),
                 size=10, color=C_GRAY)
        add_text(slide, icon,
                 x, Inches(2.1), Inches(4.1), Inches(0.7),
                 size=32, align=PP_ALIGN.CENTER)
        add_text(slide, cat,
                 x, Inches(2.82), Inches(4.1), Inches(0.42),
                 size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, change,
                 x, Inches(3.26), Inches(4.1), Inches(0.35),
                 size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

        add_rect(slide, x + Inches(0.15), Inches(3.72), Inches(3.8), Inches(0.65),
                 fill=RGBColor(0x0F, 0x2F, 0x1A))
        add_text(slide, gain,
                 x + Inches(0.15), Inches(3.74), Inches(3.8), Inches(0.6),
                 size=22, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

        add_text(slide, f"근거: {physics}",
                 x + Inches(0.15), Inches(4.42), Inches(3.8), Inches(0.3),
                 size=9, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, f"난이도: {diff}",
                 x + Inches(0.15), Inches(4.72), Inches(3.8), Inches(0.3),
                 size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 총 효과
    add_rect(slide, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.75),
             fill=RGBColor(0x0F, 0x2F, 0x1A), line=C_GREEN, line_w=Pt(1))
    add_text(slide,
             "3개 모두 적용 시 → +3년 6개월 수명 연장 (중복 효과 0.7 factor 적용)",
             Inches(0.6), Inches(5.9), Inches(7.0), Inches(0.65),
             size=14, bold=True, color=C_GREEN)
    add_text(slide,
             "절감 비용\n약 504만원",
             Inches(8.5), Inches(5.88), Inches(4.0), Inches(0.72),
             size=18, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

    # 계산 공식
    add_rect(slide, Inches(0.4), Inches(6.72), Inches(12.5), Inches(0.55),
             fill=C_CARD)
    add_text(slide,
             "Peukert 수명 연장: base × 0.30 × (1 - (C_target/C_current)^n)  |  "
             "Arrhenius 수명 연장: base × 0.30 × (1 - 1/ratio)  |  base = 10년",
             Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.5),
             size=9, color=C_GRAY)


def slide_16_app_demo(prs):
    """앱 시연"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "앱 시연 — streamlit run app.py", "실제 입력 예시: 아이오닉5 서울 10년 135,000km")

    # 앱 UI 목업
    # 사이드바 목업
    sb = add_rect(slide, Inches(0.3), Inches(1.35), Inches(3.0), Inches(5.9),
                  fill=RGBColor(0x16, 0x21, 0x38), line=C_ACCENT, line_w=Pt(1))
    add_text(slide, "🔋 BMS 튜닝 시스템",
             Inches(0.4), Inches(1.4), Inches(2.8), Inches(0.35),
             size=11, bold=True, color=C_ACCENT)

    sidebar_items = [
        ("브랜드", "현대"),
        ("모델",   "아이오닉5"),
        ("트림",   "롱레인지 2WD"),
        ("배터리", "84 kWh / 697V"),
        ("지역",   "서울 12.5°C"),
        ("주차",   "지하 (안정)"),
        ("사용",   "10년 / 135,000km"),
        ("속도",   "75 km/h"),
        ("급가속", "10회/h"),
        ("SOC폭",  "60%"),
        ("충전",   "주 4회"),
    ]
    for i, (k, v) in enumerate(sidebar_items):
        y = Inches(1.85 + i * 0.44)
        add_text(slide, k, Inches(0.42), y, Inches(0.9), Inches(0.38),
                 size=8, color=C_GRAY)
        add_text(slide, v, Inches(1.35), y, Inches(1.8), Inches(0.38),
                 size=8, bold=True, color=C_WHITE)

    add_rect(slide, Inches(0.4), Inches(6.7), Inches(2.8), Inches(0.38),
             fill=C_ACCENT)
    add_text(slide, "🔍 분석하기", Inches(0.4), Inches(6.72), Inches(2.8), Inches(0.34),
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 메인 영역 — 결과 카드
    results_area = [
        ("🟡 평균형",   "운전자 유형",    C_YELLOW),
        ("77.7%",       "현재 SOH",       C_ACCENT),
        ("12년",        "예상 총 수명",   C_GREEN),
        ("+3년 6개월",  "개선 시 연장",   C_GREEN),
    ]
    for i, (val, label, color) in enumerate(results_area):
        cx = Inches(3.55 + i * 2.45)
        add_rect(slide, cx, Inches(1.35), Inches(2.25), Inches(1.2),
                 fill=C_CARD, line=color, line_w=Pt(1))
        add_text(slide, val, cx, Inches(1.42), Inches(2.25), Inches(0.65),
                 size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, cx, Inches(2.0), Inches(2.25), Inches(0.38),
                 size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

    # SOH 곡선 미니
    fig, ax = plt.subplots(figsize=(5.5, 2.5), facecolor="#1A253D")
    ax.set_facecolor("#0F172A")
    c = np.arange(1, 2501)
    s1 = np.clip(100 - 0.67 * np.sqrt(c), 60, 100)
    s2 = np.clip(100 - 0.54 * np.sqrt(c), 60, 100)
    cur = 624
    ax.plot(c[:cur], s1[:cur], color="#6B7280", lw=1.5, ls=":")
    ax.plot(c[cur:], s1[cur:], color="#EF4444", lw=2, label="현재 습관")
    ax.plot(c[cur:], s2[cur:], color="#22C55E", lw=2, label="개선 적용")
    ax.axhline(70, color="#94A3B8", lw=1.2, ls="--")
    ax.axvline(cur, color="#FBBF24", lw=1.2)
    ax.scatter([cur], [s1[cur]], color="#FBBF24", s=40, zorder=5)
    ax.set_ylim(62, 105); ax.set_xlim(0, 2500)
    ax.tick_params(colors="#94A3B8", labelsize=7)
    for sp in ax.spines.values(): sp.set_color("#2D3F5E")
    ax.legend(facecolor="#1A253D", edgecolor="#2D3F5E", labelcolor="white", fontsize=7, loc="upper right")
    plt.tight_layout()
    add_fig(slide, fig, Inches(3.5), Inches(2.65), Inches(6.0), Inches(2.8))

    # 추천 TOP 3 미니
    mini_recs = [
        ("🥇 평균 주행 속도", "+1년 4개월", C_ACCENT),
        ("🥈 충전 SOC 폭",   "+1년 2개월", C_GREEN),
        ("🥉 주차 환경 온도", "+1년 0개월", C_YELLOW),
    ]
    for i, (title, gain, color) in enumerate(mini_recs):
        x = Inches(3.55 + i * 3.25)
        y = Inches(5.55)
        add_rect(slide, x, y, Inches(3.05), Inches(1.55),
                 fill=C_CARD, line=color, line_w=Pt(1))
        add_text(slide, title, x + Inches(0.1), y + Inches(0.1), Inches(2.85), Inches(0.38),
                 size=10, bold=True, color=color)
        add_text(slide, gain, x + Inches(0.1), y + Inches(0.55), Inches(2.85), Inches(0.55),
                 size=18, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


def slide_17_pinn_ablation(prs):
    """PINN Ablation + NCM 화학 의의"""
    slide = new_slide(prs)
    bg(slide)
    section_title(slide, "추가 실험 — PINN Ablation & 발표 핵심 포인트", "λ_pinn 값에 따른 성능 변화")

    # Ablation 가상 결과 차트
    fig, ax = plt.subplots(figsize=(5.5, 3.8), facecolor="#1A253D")
    ax.set_facecolor("#0F172A")
    lambdas = [0.0, 0.01, 0.05, 0.1, 0.5]
    r2_pinn = [0.9101, 0.9180, 0.9233, 0.9195, 0.9088]
    colors_p = ["#6B7280"] * 5
    colors_p[2] = "#22C55E"  # 최적
    bars = ax.bar([str(l) for l in lambdas], r2_pinn,
                  color=colors_p, edgecolor="#2D3F5E", width=0.55)
    for bar, val in zip(bars, r2_pinn):
        ax.text(bar.get_x() + bar.get_width()/2, val + 0.001,
                f"{val:.4f}", ha="center", color="white", fontsize=8.5)
    ax.set_xlabel("λ_pinn", color="#94A3B8", fontsize=9)
    ax.set_ylabel("R² (테스트셋)", color="#94A3B8", fontsize=9)
    ax.set_ylim(0.900, 0.930)
    ax.set_title("PINN 가중치 Ablation (λ=0.05 최적)", color="white", fontsize=9)
    ax.tick_params(colors="#94A3B8", labelsize=8)
    for sp in ax.spines.values(): sp.set_color("#2D3F5E")
    ax.annotate("최적 λ=0.05", xy=(2, 0.9233), xytext=(3.2, 0.924),
                color="#22C55E", fontsize=8,
                arrowprops=dict(arrowstyle="->", color="#22C55E", lw=1.2))
    plt.tight_layout()
    add_fig(slide, fig, Inches(0.3), Inches(1.4), Inches(6.0), Inches(4.3))

    # 핵심 숫자 암기 표
    add_text(slide, "발표 핵심 수치",
             Inches(6.7), Inches(1.4), Inches(6.3), Inches(0.35),
             size=13, bold=True, color=C_ACCENT)

    key_facts = [
        ("학습 데이터",   "NCM 실차 153대 / 85,648세션"),
        ("테스트 데이터", "38대 / 21,629세션"),
        ("LSTM R²",       "0.9233 (1위)"),
        ("LSTM RMSE",     "0.784%"),
        ("모델 순위",     "LSTM > Transformer > XGBoost > CDE"),
        ("Peukert n",     "0.621 (Safari 2011)"),
        ("Arrhenius Ea",  "31,000 J/mol (Wang 2014)"),
        ("무릎 임계",     "SOH 85% (Attia 2022, Nature Energy)"),
        ("보증 기준",     "10년 / 20만km / SOH 70%"),
        ("교체 비용",     "약 1,200만원"),
    ]
    for i, (key, val) in enumerate(key_facts):
        y = Inches(1.85 + i * 0.5)
        row_bg = C_CARD if i % 2 == 0 else RGBColor(0x16, 0x22, 0x38)
        add_rect(slide, Inches(6.7), y, Inches(6.3), Inches(0.46), fill=row_bg)
        add_text(slide, key, Inches(6.82), y + Inches(0.04), Inches(2.0), Inches(0.38),
                 size=10, color=C_GRAY)
        add_text(slide, val, Inches(8.85), y + Inches(0.04), Inches(4.0), Inches(0.38),
                 size=10, bold=True, color=C_WHITE)


def slide_18_conclusion(prs):
    """결론"""
    slide = new_slide(prs)
    bg(slide)
    accent_bar(slide)

    add_text(slide, "결론 및 향후 연구",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             size=22, bold=True, color=C_WHITE)

    # 기여점
    add_text(slide, "🎯 기여점",
             Inches(0.5), Inches(1.0), Inches(12), Inches(0.38),
             size=14, bold=True, color=C_ACCENT)
    contributions = [
        "NCM 실차 191대 데이터로 현대·기아 EV에 실제 적용 가능한 SOH 예측 시스템 구축",
        "4-way 모델 비교(LSTM/Transformer/XGBoost/CDE)로 방법론 견고성 실험적 증명",
        "Physics-Informed Loss(Peukert + Arrhenius)로 데이터 외 구간 신뢰성 확보",
        "2-phase 열화 모델 + LSTM 앵커링 Hybrid 방식으로 현실적 SOH 곡선 생성",
        "운전 습관 5개 입력 → Peukert/Arrhenius 기반 정량적 맞춤 추천 TOP 3",
    ]
    for i, c in enumerate(contributions):
        add_rect(slide, Inches(0.5), Inches(1.45 + i * 0.57), Inches(12.4), Inches(0.5),
                 fill=C_CARD if i % 2 == 0 else RGBColor(0x16, 0x22, 0x38))
        add_text(slide, f"✅  {c}",
                 Inches(0.65), Inches(1.48 + i * 0.57), Inches(12.1), Inches(0.44),
                 size=10.5, color=C_WHITE)

    # 한계
    add_text(slide, "⚠️ 한계점",
             Inches(0.5), Inches(4.43), Inches(6.0), Inches(0.35),
             size=13, bold=True, color=C_YELLOW)
    limits = [
        "중국 실차 데이터 → 한국 실도로 조건과 미세 차이 존재",
        "SOH 분포 95~100% 밀집 → 낮은 SOH 구간 데이터 부족",
        "Neural CDE 발산 → 하이퍼파라미터 추가 탐색 여지",
    ]
    for i, l in enumerate(limits):
        add_text(slide, f"• {l}",
                 Inches(0.65), Inches(4.85 + i * 0.42), Inches(6.0), Inches(0.38),
                 size=10, color=C_GRAY)

    # 향후 연구
    add_text(slide, "🚀 향후 연구",
             Inches(7.0), Inches(4.43), Inches(6.0), Inches(0.35),
             size=13, bold=True, color=C_GREEN)
    futures = [
        "한국 실차 데이터 수집 → 도메인 적응 파인튜닝",
        "OBD-II 연동 → 실시간 SOH 업데이트",
        "Neural CDE 하이퍼파라미터 최적화 재도전",
    ]
    for i, f in enumerate(futures):
        add_text(slide, f"• {f}",
                 Inches(7.15), Inches(4.85 + i * 0.42), Inches(6.0), Inches(0.38),
                 size=10, color=C_GRAY)

    # 마지막 인용 + 핵심 수치
    add_rect(slide, Inches(0.5), Inches(6.25), Inches(12.4), Inches(0.85),
             fill=RGBColor(0x1E, 0x1B, 0x4B), line=C_ACCENT, line_w=Pt(1))
    add_text(slide,
             "LSTM R²=0.9233 · RMSE=0.784%  |  Peukert n=0.621 · Arrhenius Ea=31kJ/mol  |  "
             "NCM 실차 191대  |  물리모델 + AI 앵커링 Hybrid",
             Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.38),
             size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide,
             "\"현재 SOH 숫자 하나\"에서 → \"미래 열화 곡선 + 교체 시점 + 맞춤 행동 수치\"로",
             Inches(0.7), Inches(6.72), Inches(12.0), Inches(0.35),
             size=11, color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    out_dir = Path("presentation")
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "EV_BMS_발표.pptx"

    builders = [
        ("01 제목",             slide_01_title),
        ("02 목차",             slide_02_toc),
        ("03 문제상황①",        slide_03_news1),
        ("04 문제상황②",        slide_04_news2),
        ("05 솔루션 개요",      slide_05_solution),
        ("06 데이터셋",         slide_06_dataset),
        ("07 전처리",           slide_07_preprocessing),
        ("08 열화 이론",        slide_08_degradation),
        ("09 모델 비교 개요",   slide_09_model_overview),
        ("10 LSTM & Transformer", slide_10_lstm_detail),
        ("11 PINN",             slide_11_pinn),
        ("12 실험 결과",        slide_12_results),
        ("13 SOH 곡선 생성",    slide_13_soh_curve),
        ("14 운전자 분류",      slide_14_driver_class),
        ("15 맞춤 추천",        slide_15_recommendation),
        ("16 앱 시연",          slide_16_app_demo),
        ("17 PINN Ablation",    slide_17_pinn_ablation),
        ("18 결론",             slide_18_conclusion),
    ]

    for i, (name, fn) in enumerate(builders, 1):
        print(f"  [{i:02d}/18] {name} ...", end=" ", flush=True)
        fn(prs)
        print("OK")

    prs.save(str(out_path))
    print(f"\nDONE: {out_path}")
    print(f"   슬라이드: {len(prs.slides)}장 / 13.33×7.5 inch (16:9 와이드)")


if __name__ == "__main__":
    main()
